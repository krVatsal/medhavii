"""
Quiz Generator from PPT/PDF using KimiK2

pip install PyPDF2 python-pptx groq
"""

import os
from PyPDF2 import PdfReader
from pptx import Presentation
from groq import Groq
import json

class QuizGenerator:
    def __init__(self, groq_api_key):
        """Initialize with Groq API key"""
        self.client = Groq(api_key=groq_api_key)
        
    def extract_pdf_content(self, pdf_path, page_range):
        """
        Extract text from PDF for specified page range
        
        Args:
            pdf_path: Path to PDF file
            page_range: Tuple (start_page, end_page) - 1-indexed
        """
        try:
            reader = PdfReader(pdf_path)
            total_pages = len(reader.pages)
            
            start, end = page_range
            if start < 1 or end > total_pages:
                raise ValueError(f"Page range must be between 1 and {total_pages}")
            
            content = []
            for i in range(start - 1, end):  # Convert to 0-indexed
                page = reader.pages[i]
                text = page.extract_text()
                content.append(f"--- Page {i + 1} ---\n{text}")
            
            return "\n\n".join(content)
        
        except Exception as e:
            raise Exception(f"Error extracting PDF: {str(e)}")
    
    def extract_ppt_content(self, ppt_path, slide_range):
        """
        Extract text from PowerPoint for specified slide range
        
        Args:
            ppt_path: Path to PPT/PPTX file
            slide_range: Tuple (start_slide, end_slide) - 1-indexed
        """
        try:
            prs = Presentation(ppt_path)
            total_slides = len(prs.slides)
            
            start, end = slide_range
            if start < 1 or end > total_slides:
                raise ValueError(f"Slide range must be between 1 and {total_slides}")
            
            content = []
            for i in range(start - 1, end):  # Convert to 0-indexed
                slide = prs.slides[i]
                slide_text = []
                
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        slide_text.append(shape.text)
                
                content.append(f"--- Slide {i + 1} ---\n" + "\n".join(slide_text))
            
            return "\n\n".join(content)
        
        except Exception as e:
            raise Exception(f"Error extracting PPT: {str(e)}")
    
    def generate_quiz(self, content, num_questions=5, difficulty="medium"):
        """
        Generate quiz using Groq LLM
        
        Args:
            content: Extracted text content
            num_questions: Number of questions to generate
            difficulty: easy, medium, or hard
        """
        prompt = f"""Based on the following content, generate {num_questions} multiple-choice quiz questions.

Difficulty level: {difficulty}

Content:
{content}

Generate a quiz in JSON format with the following structure:
{{
    "quiz": [
        {{
            "question": "Question text here?",
            "options": ["A) Option 1", "B) Option 2", "C) Option 3", "D) Option 4"],
            "correct_answer": "A",
            "explanation": "Brief explanation of the correct answer"
        }}
    ]
}}

Make sure the questions are relevant, clear, and test understanding of the key concepts.
Return ONLY the JSON, no additional text."""

        try:
            chat_completion = self.client.chat.completions.create(
                messages=[
                    {
                        "role": "system",
                        "content": "You are an expert quiz creator. Generate high-quality multiple-choice questions based on provided content."
                    },
                    {
                        "role": "user",
                        "content": prompt
                    }
                ],
                model="moonshotai/kimi-k2-instruct-0905",
                temperature=0.7,
                max_tokens=8192
            )
            
            response_text = chat_completion.choices[0].message.content
            
            
            start_idx = response_text.find('{')
            end_idx = response_text.rfind('}') + 1
            json_str = response_text[start_idx:end_idx]
            
            quiz_data = json.loads(json_str)
            return quiz_data
        
        except Exception as e:
            raise Exception(f"Error generating quiz: {str(e)}")
    
    def display_quiz(self, quiz_data):
        """Display quiz in a readable format"""
        print("\n" + "="*60)
        print("GENERATED QUIZ")
        print("="*60 + "\n")
        
        for i, q in enumerate(quiz_data['quiz'], 1):
            print(f"Question {i}: {q['question']}")
            for option in q['options']:
                print(f"  {option}")
            print(f"\nCorrect Answer: {q['correct_answer']}")
            print(f"Explanation: {q['explanation']}")
            print("\n" + "-"*60 + "\n")


def main(path,range,q_no,difficult):
 
    
    # Set your Groq API key
    GROQ_API_KEY = "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"  
    
    # Initialize generator
    generator = QuizGenerator(GROQ_API_KEY)
    
 
    pdf_path = path#"/content/Vinay K. Ingle, John G. Proakis - Digital Signal Processing Using MATLAB, 3rd Edition  -Cengage Learning (2011).pdf"
    pdf_page_range = (range[0],range[1])  
    
    try:
        content = generator.extract_pdf_content(pdf_path, pdf_page_range)
        print(f"Extracted {len(content)} characters from PDF")
        
        quiz = generator.generate_quiz(content, num_questions=q_no, difficulty=difficult)
        generator.display_quiz(quiz)
        
    except Exception as e:
        print(f"Error: {e}")

    
    ppt_path = path
    ppt_slide_range = (range[0], range[1])  
    
    try:
        content = generator.extract_ppt_content(ppt_path, ppt_slide_range)
        print(f"Extracted {len(content)} characters from PowerPoint")
        
        quiz = generator.generate_quiz(content, num_questions=q_no, difficulty=difficult)
        generator.display_quiz(quiz)
        
    except Exception as e:
        print(f"Error: {e}")
    
    # Save quiz to JSON file
    with open('quiz_output.json', 'w') as f:
        json.dump(quiz, f, indent=2)



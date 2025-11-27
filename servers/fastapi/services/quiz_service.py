"""
Quiz Generator Service from Presentation Slides using Groq LLM
"""
import os
import json
from typing import List, Tuple, Optional, Literal
from groq import Groq
from pptx import Presentation


class QuizGenerationService:
    """Service to generate quiz questions from presentation slides"""
    
    def __init__(self):
        """Initialize with Groq API key from environment"""
        self.groq_api_key = os.getenv("GROQ_API_KEY", "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")
        if not self.groq_api_key:
            raise ValueError("GROQ_API_KEY environment variable is not set")
        self.client = Groq(api_key=self.groq_api_key)
    
    def extract_ppt_content(self, ppt_path: str, slide_range: Tuple[int, int]) -> str:
        """
        Extract text from PowerPoint for specified slide range
        
        Args:
            ppt_path: Path to PPT/PPTX file
            slide_range: Tuple (start_slide, end_slide) - 1-indexed
            
        Returns:
            Extracted text content from slides
        """
        try:
            prs = Presentation(ppt_path)
            total_slides = len(prs.slides)
            
            start, end = slide_range
            if start < 1 or end > total_slides:
                raise ValueError(f"Slide range must be between 1 and {total_slides}")
            
            content = []
            for i in range(start - 1, end):  # Convert to 0-indexed
                slide = prs.slides[i]
                slide_text = []
                
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text.strip():
                        slide_text.append(shape.text)
                
                if slide_text:  # Only add if there's actual text
                    content.append(f"--- Slide {i + 1} ---\n" + "\n".join(slide_text))
            
            return "\n\n".join(content)
        
        except Exception as e:
            raise Exception(f"Error extracting PPT content: {str(e)}")
    
    def extract_slide_text_from_data(self, slides_data: List[dict]) -> str:
        """
        Extract text from slides data (JSON format from database)
        
        Args:
            slides_data: List of slide dictionaries with content
            
        Returns:
            Extracted text content from slides
        """
        try:
            content = []
            for i, slide in enumerate(slides_data, 1):
                slide_text = []
                
                # Extract text from various fields
                if "title" in slide and slide["title"]:
                    slide_text.append(f"Title: {slide['title']}")
                
                if "subtitle" in slide and slide["subtitle"]:
                    slide_text.append(f"Subtitle: {slide['subtitle']}")
                
                if "content" in slide and slide["content"]:
                    if isinstance(slide["content"], list):
                        slide_text.extend(slide["content"])
                    else:
                        slide_text.append(str(slide["content"]))
                
                if "description" in slide and slide["description"]:
                    slide_text.append(f"Description: {slide['description']}")
                
                if "bullets" in slide and slide["bullets"]:
                    slide_text.append("Bullets:")
                    if isinstance(slide["bullets"], list):
                        slide_text.extend([f"• {bullet}" for bullet in slide["bullets"]])
                
                if slide_text:
                    content.append(f"--- Slide {i} ---\n" + "\n".join(slide_text))
            
            return "\n\n".join(content)
        
        except Exception as e:
            raise Exception(f"Error extracting slide data: {str(e)}")
    
    async def generate_quiz(
        self,
        content: str,
        num_questions: int = 5,
        difficulty: Literal["easy", "medium", "hard"] = "medium"
    ) -> dict:
        """
        Generate quiz using Groq LLM (async version)
        
        Args:
            content: Extracted text content
            num_questions: Number of questions to generate (1-10)
            difficulty: Difficulty level (easy, medium, or hard)
            
        Returns:
            Dictionary containing quiz questions
        """
        # Validate inputs
        if num_questions < 1 or num_questions > 10:
            raise ValueError("Number of questions must be between 1 and 10")
        
        if difficulty not in ["easy", "medium", "hard"]:
            raise ValueError("Difficulty must be 'easy', 'medium', or 'hard'")
        
        prompt = f"""Based on the following content, generate {num_questions} multiple-choice quiz questions.

Difficulty level: {difficulty}

Content:
{content}

Generate a quiz in JSON format with the following structure:
{{
    "quiz": [
        {{
            "question": "Question text here?",
            "options": ["A) Option 1", "B) Option 2", "C) Option 3", "D) Option 4"],
            "correct_answer": "A",
            "explanation": "Brief explanation of the correct answer"
        }}
    ]
}}

Make sure the questions are relevant, clear, and test understanding of the key concepts from the content.
For {difficulty} difficulty:
- easy: Basic recall and understanding questions
- medium: Application and analysis questions
- hard: Synthesis and evaluation questions

Return ONLY the JSON, no additional text or markdown formatting."""

        try:
            chat_completion = self.client.chat.completions.create(
                messages=[
                    {
                        "role": "system",
                        "content": "You are an expert quiz creator. Generate high-quality multiple-choice questions based on provided content. Return only valid JSON without any markdown formatting or additional text."
                    },
                    {
                        "role": "user",
                        "content": prompt
                    }
                ],
                model="moonshotai/kimi-k2-instruct-0905",
                temperature=0.7,
                max_tokens=2048
            )
            
            response_text = chat_completion.choices[0].message.content
            
            # Clean the response - remove markdown code blocks if present
            response_text = response_text.strip()
            if response_text.startswith("```json"):
                response_text = response_text[7:]
            if response_text.startswith("```"):
                response_text = response_text[3:]
            if response_text.endswith("```"):
                response_text = response_text[:-3]
            response_text = response_text.strip()
            
            # Extract JSON from response
            start_idx = response_text.find('{')
            end_idx = response_text.rfind('}') + 1
            
            if start_idx == -1 or end_idx == 0:
                raise ValueError("No valid JSON found in response")
            
            json_str = response_text[start_idx:end_idx]
            quiz_data = json.loads(json_str)
            
            # Validate quiz structure
            if "quiz" not in quiz_data or not isinstance(quiz_data["quiz"], list):
                raise ValueError("Invalid quiz structure in response")
            
            return quiz_data
        
        except json.JSONDecodeError as e:
            raise Exception(f"Error parsing quiz JSON: {str(e)}")
        except Exception as e:
            raise Exception(f"Error generating quiz: {str(e)}")
    
    def generate_quiz_sync(
        self,
        content: str,
        num_questions: int = 5,
        difficulty: Literal["easy", "medium", "hard"] = "medium"
    ) -> dict:
        """
        Generate quiz using Groq LLM (sync version)
        
        Same as generate_quiz but synchronous for compatibility
        """
        # Validate inputs
        if num_questions < 1 or num_questions > 10:
            raise ValueError("Number of questions must be between 1 and 10")
        
        if difficulty not in ["easy", "medium", "hard"]:
            raise ValueError("Difficulty must be 'easy', 'medium', or 'hard'")
        
        prompt = f"""Based on the following content, generate {num_questions} multiple-choice quiz questions.

Difficulty level: {difficulty}

Content:
{content}

Generate a quiz in JSON format with the following structure:
{{
    "quiz": [
        {{
            "question": "Question text here?",
            "options": ["A) Option 1", "B) Option 2", "C) Option 3", "D) Option 4"],
            "correct_answer": "A",
            "explanation": "Brief explanation of the correct answer"
        }}
    ]
}}

Make sure the questions are relevant, clear, and test understanding of the key concepts.
Return ONLY the JSON, no additional text or markdown formatting."""

        try:
            chat_completion = self.client.chat.completions.create(
                messages=[
                    {
                        "role": "system",
                        "content": "You are an expert quiz creator. Generate high-quality multiple-choice questions based on provided content. Return only valid JSON."
                    },
                    {
                        "role": "user",
                        "content": prompt
                    }
                ],
                model="moonshotai/kimi-k2-instruct-0905",
                temperature=0.7,
                max_tokens=2048
            )
            
            response_text = chat_completion.choices[0].message.content
            
            # Clean and extract JSON
            response_text = response_text.strip()
            if response_text.startswith("```json"):
                response_text = response_text[7:]
            if response_text.startswith("```"):
                response_text = response_text[3:]
            if response_text.endswith("```"):
                response_text = response_text[:-3]
            response_text = response_text.strip()
            
            start_idx = response_text.find('{')
            end_idx = response_text.rfind('}') + 1
            
            if start_idx == -1 or end_idx == 0:
                raise ValueError("No valid JSON found in response")
            
            json_str = response_text[start_idx:end_idx]
            quiz_data = json.loads(json_str)
            
            if "quiz" not in quiz_data or not isinstance(quiz_data["quiz"], list):
                raise ValueError("Invalid quiz structure in response")
            
            return quiz_data
        
        except json.JSONDecodeError as e:
            raise Exception(f"Error parsing quiz JSON: {str(e)}")
        except Exception as e:
            raise Exception(f"Error generating quiz: {str(e)}")


# Global service instance
quiz_service = QuizGenerationService()
